# -*- coding: utf-8 -*-
"""Editable PowerPoint version of the data-driven graphical abstract (v3).

Rebuilds the v3 layout with NATIVE PowerPoint shapes/text (fully editable in
PowerPoint -- no embedded raster image), reading the SAME backbone-agnostic
JSON data as make_data_driven_GA_v3.py so the numbers stay data-grounded.

Outputs: figures/graphical_abstract_v3_editable.pptx
Does NOT touch data_driven_GA_v3.png / .pdf or graphical_abstract_v2.pptx.
"""
import json
from pathlib import Path
import numpy as np
from scipy import stats

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# adjust path: backbone-agnostic summary JSON lives in the analysis workspace (NOT in repo)
WS = Path("ADJUST_PATH/workspace")
OUT = Path("./graphical_abstract/graphical_abstract_v3_editable.pptx")  # adjust to repo graphical_abstract/ dir

# ---------- Extract data (mirror make_data_driven_GA_v3.py) ----------
summary = json.loads((WS / "p2b_backbone_agnostic_summary.json").read_text())
data = {}
for bb, d in summary["by_backbone"].items():
    data[bb] = {
        "lesion": (d["lesion_evidence_mean"]["mean"], d["lesion_evidence_mean"]["sd"]),
        "fp":     (d["fp_evidence_mean"]["mean"],     d["fp_evidence_mean"]["sd"]),
        "benign": (d["benign_evidence_mean"]["mean"], d["benign_evidence_mean"]["sd"]),
        "n": d["n"],
    }

ref_dir = WS / "p2b_archive" / "workspace" / "vast_5fold" / "results"
ref_lesion, ref_fp, ref_benign = [], [], []
for f in sorted(ref_dir.glob("p2b_fold*_seed*_picai.json")):
    d = json.loads(f.read_text())
    scen = d.get("scenarios", {}).get("ideal", {})
    fp_block = d.get("fp_evidence_ideal", {})
    fm = fp_block.get("fp_evidence", {}).get("mean")
    lm = scen.get("lesion_evidence_mean")
    bm = scen.get("benign_evidence_mean")
    if all(v is not None for v in [fm, lm, bm]):
        ref_lesion.append(lm); ref_fp.append(fm); ref_benign.append(bm)
data["migf_nnunet_A2"] = {
    "lesion": (float(np.mean(ref_lesion)), float(np.std(ref_lesion, ddof=1))),
    "fp":     (float(np.mean(ref_fp)),     float(np.std(ref_fp,     ddof=1))),
    "benign": (float(np.mean(ref_benign)), float(np.std(ref_benign, ddof=1))),
    "n": len(ref_lesion),
    "raw_lesion": ref_lesion, "raw_fp": ref_fp, "raw_benign": ref_benign,
}

display_order = ["bare_nnunet", "bare_unet", "bare_mamba", "migf_mamba_A2", "migf_nnunet_A2"]
display_label = {
    "bare_nnunet": "bare nnU-Net", "bare_unet": "bare U-Net", "bare_mamba": "bare Mamba",
    "migf_mamba_A2": "MIGF-Mamba A2", "migf_nnunet_A2": "MIGF-nnUNet A2",
}

# normalized lesion position per backbone (benign=0, FP=1) -- drives circle x-positions
lesion_norm = {}
for bb in display_order:
    d = data[bb]; lm = d["lesion"][0]; fm = d["fp"][0]; bm = d["benign"][0]
    span = fm - bm
    lesion_norm[bb] = (lm - bm) / span if span else 0.0

# Effect sizes shown in the pptx are the manuscript-reported values
# (Cohen's d 1.10 lesion-vs-benign; FP/benign evidence ratio 2.38x; 35/35).
for bb in display_order:
    print("  %-16s lesion_norm=%.2f  n=%d" % (bb, lesion_norm[bb], data[bb]["n"]))

# ---------- Colors ----------
TEAL_DARK = RGBColor(0x0F, 0x4C, 0x5C)
RED_LESION = RGBColor(0xB9, 0x1C, 0x1C)
ORANGE_FP = RGBColor(0xEA, 0x58, 0x0C)
GRAY_BENIGN = RGBColor(0x6B, 0x72, 0x80)
GOLD = RGBColor(0xC8, 0x8B, 0x1E)
DARK_TEXT = RGBColor(0x37, 0x41, 0x51)
BOX_BG = RGBColor(0xF9, 0xFA, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- Slide ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def add_text(left, top, width, height, text, *, font_size=12, bold=False, italic=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font_name; r.font.size = Pt(font_size)
        r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb


def add_rect(left, top, width, height, *, fill=None, line=None, line_w=1.0, rounded=False):
    shape_t = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_t, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_circle(cx, cy, d, *, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.color.rgb = WHITE; s.line.width = Pt(1.5)
    s.shadow.inherit = False
    return s


def add_dashed(x0, y0, x1, y1, *, color=RGBColor(0xCB, 0xD2, 0xD8)):
    from pptx.oxml.ns import qn
    ln_shape = slide.shapes.add_connector(1, Inches(x0), Inches(y0), Inches(x1), Inches(y1))
    ln_shape.line.color.rgb = color; ln_shape.line.width = Pt(0.75)
    ln = ln_shape.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return ln_shape


# ---------- Title bar ----------
add_rect(0, 0, 13.333, 0.85, fill=TEAL_DARK)
add_text(0.4, 0.13, 12.6, 0.6,
         "False positives are contrast-matched to cancer — a data-level imaging property across 5 architectures",
         font_size=20, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ---------- Left panel geometry ----------
PLOT_X0 = 3.35   # normalized 0 (benign mean)
PLOT_X1 = 7.45   # normalized 1.0 (FP mean)
def x_of(norm):
    return PLOT_X0 + norm * (PLOT_X1 - PLOT_X0)

ROW_CY = [1.98, 2.60, 3.22, 3.84, 4.46]
DOT_D = 0.30

# Column headers
add_text(x_of(0.0) - 0.7, 1.42, 1.4, 0.3, "Benign", font_size=11, bold=True,
         color=GRAY_BENIGN, align=PP_ALIGN.CENTER)
add_text(x_of(0.5) - 0.9, 1.42, 1.8, 0.3, "Lesion (position)", font_size=11, bold=True,
         color=RED_LESION, align=PP_ALIGN.CENTER)
add_text(x_of(1.0) - 0.8, 1.42, 1.6, 0.3, "False positive", font_size=11, bold=True,
         color=ORANGE_FP, align=PP_ALIGN.CENTER)

# Rows
for cy, bb in zip(ROW_CY, display_order):
    ln_norm = lesion_norm[bb]
    add_dashed(x_of(-0.02), cy, x_of(1.05), cy)
    add_circle(x_of(0.0), cy, DOT_D, color=GRAY_BENIGN)
    add_circle(x_of(ln_norm), cy, DOT_D, color=RED_LESION)
    add_circle(x_of(1.0), cy, DOT_D, color=ORANGE_FP)
    # lesion value label above lesion dot
    add_text(x_of(ln_norm) - 0.45, cy - 0.50, 0.9, 0.28, f"{ln_norm:.2f}",
             font_size=9, bold=True, color=RED_LESION, align=PP_ALIGN.CENTER)
    # backbone label (right-aligned, left of benign)
    add_text(0.35, cy - 0.16, 2.75, 0.34, display_label[bb], font_size=11, bold=True,
             color=DARK_TEXT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # n label (right of FP)
    add_text(x_of(1.0) + 0.20, cy - 0.14, 0.9, 0.3, f"n={data[bb]['n']}", font_size=9,
             italic=True, color=GRAY_BENIGN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# x tick labels
add_text(x_of(0.0) - 0.8, 4.78, 1.6, 0.5, "0\n(benign mean)", font_size=8.5,
         color=GRAY_BENIGN, align=PP_ALIGN.CENTER)
add_text(x_of(0.5) - 0.4, 4.78, 0.8, 0.3, "0.50", font_size=8.5, color=GRAY_BENIGN, align=PP_ALIGN.CENTER)
add_text(x_of(1.0) - 0.8, 4.78, 1.6, 0.5, "1.0\n(FP mean)", font_size=8.5,
         color=GRAY_BENIGN, align=PP_ALIGN.CENTER)
# x-axis title
add_text(2.6, 5.30, 5.7, 0.3,
         "Normalized evidence score (per backbone: benign mean = 0, FP mean = 1)",
         font_size=9.5, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ---------- Right evidence panel ----------
RB_L, RB_T, RB_W, RB_H = 8.85, 1.42, 4.10, 3.85
add_rect(RB_L, RB_T, RB_W, RB_H, fill=BOX_BG, line=GOLD, line_w=1.6, rounded=True)
rcx = RB_L  # use full-width centered text boxes
add_text(RB_L, RB_T + 0.12, RB_W, 0.7, "35 / 35", font_size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 0.86, RB_W, 0.28, "observations reproduce", font_size=11, color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 1.12, RB_W, 0.28, "benign < lesion < FP", font_size=11, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 1.38, RB_W, 0.40,
         "20 backbone-fold + 15 reference\nobservations = 35",
         font_size=8.5, italic=True, color=GRAY_BENIGN, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 1.86, RB_W, 0.26, "Statistical evidence (main text)", font_size=10.5, bold=True,
         color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 2.16, RB_W, 0.24, "Cohen's d (lesion vs benign) = 1.10", font_size=9,
         color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 2.40, RB_W, 0.24, "FP / benign evidence ratio = 2.38×", font_size=9,
         color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 2.64, RB_W, 0.24, "Both directions reproduce 35 / 35", font_size=9,
         color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(RB_L, RB_T + 2.98, RB_W, 0.26, "Why are false positives cancer-like?", font_size=9.5, bold=True,
         color=TEAL_DARK, align=PP_ALIGN.CENTER)
why = ("False positives are the backbone's above-threshold errors — regions pre-selected "
       "for strong cancer-like imaging features, so their raw T2W/ADC contrast matches "
       "true cancer far more than benign tissue (a data-level imaging property).")
add_text(RB_L + 0.18, RB_T + 3.24, RB_W - 0.36, 0.6, why, font_size=8, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ---------- Bottom caption ----------
caption = ("Each row is one backbone architecture. The three circles mark the mean evidence-head score for "
           "benign tissue (gray), the true lesion (red), and the model's own false positives (orange), per-backbone "
           "normalized so benign = 0 and FP = 1. In all five architectures both the lesion and the false positives "
           "separate sharply from benign tissue — i.e., false positives are contrast-matched to confirmed cancer.")
add_text(0.4, 5.70, 12.5, 1.0, caption, font_size=9, color=DARK_TEXT, align=PP_ALIGN.LEFT)

# ---------- Footer ----------
add_rect(0, 7.05, 13.333, 0.45, fill=TEAL_DARK)
add_text(0.4, 7.12, 12.5, 0.3,
         "Shu Y et al. — Central South University — European Radiology submission, 2026",
         font_size=10, italic=True, color=WHITE)

prs.save(OUT)
print("\nSaved: %s (%.1f KB)" % (OUT, OUT.stat().st_size / 1024))
